from pptx.enum.shapes import MSO_SHAPE_TYPE

def iter_shapes(shape):
    yield shape
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            yield from iter_shapes(s)

def collect_run_targets(prs, include_notes=True):
    """
    Returns list of (run_object, original_text).
    We replace run.text to preserve formatting.
    """
    targets = []
    for slide in prs.slides:
        # shapes
        for shape in slide.shapes:
            for s in iter_shapes(shape):
                # text frames
                if hasattr(s, "has_text_frame") and s.has_text_frame:
                    tf = s.text_frame
                    for p in tf.paragraphs:
                        for run in p.runs:
                            if run.text and run.text.strip():
                                targets.append((run, run.text))

                # tables
                if hasattr(s, "has_table") and s.has_table:
                    table = s.table
                    for row in table.rows:
                        for cell in row.cells:
                            tf = cell.text_frame
                            for p in tf.paragraphs:
                                for run in p.runs:
                                    if run.text and run.text.strip():
                                        targets.append((run, run.text))

        # notes
        if include_notes and slide.has_notes_slide and slide.notes_slide:
            ntf = slide.notes_slide.notes_text_frame
            if ntf:
                for p in ntf.paragraphs:
                    for run in p.runs:
                        if run.text and run.text.strip():
                            targets.append((run, run.text))

    return targets