from pptx.enum.shapes import MSO_SHAPE_TYPE


def iter_shapes(shape):
    """Recursively yield shapes (supports GROUP shapes)."""
    yield shape
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            yield from iter_shapes(s)


def collect_run_targets(prs, include_notes=True):
    """
    Returns list of (run_object, original_text)
    - Best formatting preservation (run-level)
    """
    targets = []

    for slide in prs.slides:
        for shape in slide.shapes:
            for s in iter_shapes(shape):

                # Text frames
                if hasattr(s, "has_text_frame") and s.has_text_frame:
                    tf = s.text_frame
                    for p in tf.paragraphs:
                        for run in p.runs:
                            if run.text is not None and run.text.strip() != "":
                                targets.append((run, run.text))

                # Tables
                if hasattr(s, "has_table") and s.has_table:
                    table = s.table
                    for row in table.rows:
                        for cell in row.cells:
                            tf = cell.text_frame
                            for p in tf.paragraphs:
                                for run in p.runs:
                                    if run.text is not None and run.text.strip() != "":
                                        targets.append((run, run.text))

        # Notes
        if include_notes and slide.has_notes_slide and slide.notes_slide:
            ntf = slide.notes_slide.notes_text_frame
            if ntf:
                for p in ntf.paragraphs:
                    for run in p.runs:
                        if run.text is not None and run.text.strip() != "":
                            targets.append((run, run.text))

    return targets


def collect_paragraph_targets(prs, include_notes=True):
    """
    Returns list of (runs_list, combined_text)
    - Better translation completeness when words are split across runs
    - Slightly more formatting risk (we keep first run formatting and clear rest)
    """
    targets = []

    for slide in prs.slides:
        for shape in slide.shapes:
            for s in iter_shapes(shape):
                if hasattr(s, "has_text_frame") and s.has_text_frame:
                    tf = s.text_frame
                    for p in tf.paragraphs:
                        runs = list(p.runs)
                        if not runs:
                            continue
                        combined = "".join(r.text or "" for r in runs)
                        if combined.strip():
                            targets.append((runs, combined))

                if hasattr(s, "has_table") and s.has_table:
                    table = s.table
                    for row in table.rows:
                        for cell in row.cells:
                            tf = cell.text_frame
                            for p in tf.paragraphs:
                                runs = list(p.runs)
                                if not runs:
                                    continue
                                combined = "".join(r.text or "" for r in runs)
                                if combined.strip():
                                    targets.append((runs, combined))

        if include_notes and slide.has_notes_slide and slide.notes_slide:
            ntf = slide.notes_slide.notes_text_frame
            if ntf:
                for p in ntf.paragraphs:
                    runs = list(p.runs)
                    if not runs:
                        continue
                    combined = "".join(r.text or "" for r in runs)
                    if combined.strip():
                        targets.append((runs, combined))

    return targets
